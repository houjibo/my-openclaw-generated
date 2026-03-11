#!/usr/bin/env python3
"""
PowerPoint 创建工具
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import sys
import json

def create_presentation(content):
    """创建 PowerPoint 演示文稿"""

    try:
        prs = Presentation()

        # 标题幻灯片
        if 'title' in content:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = content['title']
            subtitle.text = content.get('subtitle', '')

        # 内容幻灯片
        if 'slides' in content:
            for slide_data in content['slides']:
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                # 标题
                if 'title' in slide_data:
                    slide.shapes.title.text = slide_data['title']

                # 内容（如果有）
                if 'content' in slide_data:
                    text_box = slide.shapes.add_textbox(
                        Inches(1), Inches(1.5), Inches(8), Inches(4)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = slide_data['content']

        # 表格幻灯片
        if 'tables' in content:
            for table_data in content['tables']:
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                if 'title' in table_data:
                    slide.shapes.title.text = table_data['title']

                rows = len(table_data['data'])
                cols = len(table_data['data'][0])

                # 添加表格
                table = slide.shapes.add_table(
                    rows, cols,
                    Inches(0.5), Inches(1.5),
                    Inches(9), Inches(4)
                ).table

                for i, row_data in enumerate(table_data['data']):
                    for j, cell_data in enumerate(row_data):
                        table.rows[i].cells[j].text = cell_data

        output_path = content.get('output_path', '/tmp/output.pptx')
        prs.save(output_path)
        return {'success': True, 'output_path': output_path}

    except Exception as e:
        return {'error': str(e)}

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({
            'error': 'Usage: pptx_creator.py <json_content> [output_path]'
        }, ensure_ascii=False))
        sys.exit(1)

    json_content = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else '/tmp/output.pptx'

    try:
        content = json.loads(json_content)
        content['output_path'] = output_path
        result = create_presentation(content)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except json.JSONDecodeError as e:
        print(json.dumps({'error': f'Invalid JSON: {str(e)}'}, ensure_ascii=False))
